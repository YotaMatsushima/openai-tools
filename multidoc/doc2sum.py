from flask import Flask, request, render_template, session
import os
import re
from werkzeug.utils import secure_filename
from pptx import Presentation
from docx import Document
from youtube_transcript_api import YouTubeTranscriptApi
import pdfplumber
import openai
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter


# Load environment variables
openai_key = os.getenv("OPENAI_KEY")
openai_organization = os.getenv("OPENAI_ORGANIZATION")

# OpenAI API configurations
openai.api_key = openai_key
openai.organization = openai_organization

app = Flask(__name__)
app.secret_key = os.urandom(24)  # sessionを使うためには秘密鍵が必要です。


def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = []

    for slide in presentation.slides:
        # Extract text from slide content
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

        # Extract text from slide notes
        if slide.has_notes_slide:
            for note_shape in slide.notes_slide.notes_text_frame.text:
                if hasattr(note_shape, "text"):
                    text.append(note_shape.text)

    return ' '.join(text)

def extract_text_from_docx(file):
    doc = Document(file)
    text = " ".join([p.text for p in doc.paragraphs])
    return text

def extract_text_from_pdf(file):
    with pdfplumber.open(file) as pdf:
        text = " ".join(page.extract_text() for page in pdf.pages)
    return text

def text_video(youtube_url):
    """Function to get and process the subtitles, and return the summaries."""

    # Extract video id from url
    video_id = re.findall(r'(?:v=|\/)([0-9A-Za-z_-]{11}).*', youtube_url)[0]

    print("Loading subtitles...")
    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        transcript = transcript_list.find_transcript(["ja"])
        subtitles = transcript.fetch()
    except NoTranscriptFound:
        print("Error: No subtitles found in the specified language.")
        return

    # Collect subtitles without timestamps
    text = ""
    for subtitle in subtitles:
        text += subtitle['text'] + " "
    return text

def summarize_file(file):
    """Function to extract text from file and return the summaries."""

    # Check file extension
    file_extension = os.path.splitext(file.filename)[1]
    if file_extension == '.pptx':
        text = extract_text_from_pptx(file)
    elif file_extension == '.docx':
        text = extract_text_from_docx(file)
    elif file_extension == '.pdf':
        text = extract_text_from_pdf(file)
    else:
        print("Unsupported file type. Please upload a .pptx, .docx or .pdf file.")
        return

    return text

def summaryText(text,prompt):
    chunks = []
    start = 0
    text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
        separator = " ",
        chunk_size=800,
        chunk_overlap=0
        )

    texts=text_splitter.split_text(text)

    summaries = []
    # Generate a summary for each chunk
    for chunk in texts:
        print("####chunktext:"+chunk)
        messages = [
            {"role": "system", "content": "You are a helpful assistant that summarizes text."},
            {"role": "user", "content": prompt + ": " + chunk}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=messages,
        )
        summary = response['choices'][0]['message']['content']
        summaries.append(summary)
        # Add the assistant's reply to the conversation
    session['messages'].append({"role": "assistant", "content":  "\n\n".join(summaries)})
    return "\n\n".join(summaries)


def aiconversation(text,prompt):

    if text=="":
        return "error!!!!!!"
    if prompt=="":
        prompt=os.getenv("DEFAULT_PROMPT")

    print("#####"+text+"######")
    print("#####prompt:"+prompt+"######")

   # Start a new conversation if there is no existing conversation

    if 'messages' not in session:
        session['messages'] = [
            {"role": "system", "content": "You are a helpful assistant that summarizes text."},
        ]
        return summaryText(text,prompt)

    session['messages'].append({"role": "user", "content": prompt})
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=session['messages'],
    )

    # Add the assistant's reply to the conversation
    session['messages'].append({"role": "assistant", "content": response['choices'][0]['message']['content']})

    summary = response['choices'][0]['message']['content']
    print("###output:"+summary+"####")
    return summary


@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':

        prompt = request.form.get('prompt', '')  # Get prompt from form
        text=request.form.get('youtube_url', '')  # Get prompt from form
        summary=""
        if not text=="":
            print("text_video/n")
            summary = text_video(text)
        else:
            if 'file' not in request.files:
                return 'No file part in the request.'
            file = request.files['file']
            print("file conversation start!")
            if file.filename == '':
                return 'No selected file.'
            if file:
                file.seek(0)  # Ensure we're at the start of the file.
                summary = summarize_file(file)

        return aiconversation(summary,prompt)

    return render_template('upload.html')

@app.route('/reset', methods=['POST'])
def reset_session():
    print("reset session!!")
    session.pop('messages', None)  # Remove 'messages' from the session
    return ""


if __name__ == '__main__':
    app.run(port=5000, debug=True)
